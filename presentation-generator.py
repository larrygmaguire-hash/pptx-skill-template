#!/usr/bin/env python3
"""
Data-driven PowerPoint generator for branded presentations.

This script generates slides dynamically from a structured outline.
Customise the OUTLINE dict or pass your own outline to create_presentation().

Requirements:
    pip install python-pptx lxml

Usage:
    1. Update configuration section (template path, colours, font, layouts)
    2. Modify OUTLINE with your presentation content
    3. Run: python presentation-generator.py

The script will expand sections and slides based on your outline structure.
"""

import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


# =============================================================================
# CONFIGURATION — CUSTOMISE THESE FOR YOUR TEMPLATE
# =============================================================================

# Path to your PowerPoint template (.potx or .pptx)
TEMPLATE_PATH = Path("/path/to/your-template.potx")

# Output path for generated presentations
OUTPUT_PATH = Path("./output.pptx")

# Your brand accent colour (RGB values)
ACCENT_COLOUR = RGBColor(0, 0, 0)  # Replace with your RGB values
ACCENT_HEX = "000000"  # Same colour as hex (without #)

# Your font family
FONT_FAMILY = "Arial"  # Replace with your font

# Layout index mapping — MUST match your template
# Run extraction script to get these values
LAYOUTS = {
    "title": 0,
    "menu": 1,           # Agenda/contents
    "section": 2,        # Section dividers
    "section_pale": 3,   # Section dividers (pale/light background)
    "about": 4,          # Fixed about slide
    "content_white": 5,  # Body content (white background)
    "content_pale": 6,   # Body content (pale background)
    "quote": 7,          # Quotes
    "cta": 8,            # Call-to-action (fixed)
    "thank_you": 9,      # Closing slide
}

# Layouts with fixed content (animations pre-set in template)
FIXED_LAYOUTS = ["about", "cta"]

# Body placeholder index (check your template)
BODY_PLACEHOLDER_IDX = 13


# =============================================================================
# TEMPLATE CONVERSION
# =============================================================================

def convert_potx_to_pptx(potx_path: Path, pptx_path: Path):
    """Convert .potx template to .pptx by updating content type."""
    with zipfile.ZipFile(potx_path, "r") as zin:
        with zipfile.ZipFile(pptx_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    data = data.replace(
                        b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                        b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
                    )
                zout.writestr(item, data)


# =============================================================================
# TEXT FORMATTING
# =============================================================================

def set_text_simple(placeholder, text: str, is_light_bg: bool = False, font_size: int = None):
    """
    Set text with typography and 1.2 line spacing (no bullets).
    Use for titles and subtitles.
    """
    placeholder.text = text
    for paragraph in placeholder.text_frame.paragraphs:
        paragraph.line_spacing = 1.2
        for run in paragraph.runs:
            run.font.name = FONT_FAMILY
            if font_size:
                run.font.size = Pt(font_size)
            if is_light_bg:
                run.font.color.rgb = ACCENT_COLOUR


def set_body_with_bullets(placeholder, intro_text: str, bullets: list, is_light_bg: bool = False):
    """
    Set body content with intro paragraph, blank line, then arrow bullets.

    Structure:
        - Intro paragraph (1-2 sentences)
        - Blank line for spacing
        - Arrow bullet points
    """
    placeholder.text = ""
    txBody = placeholder._element.find(qn('p:txBody'))
    if txBody is None:
        return

    # Remove existing paragraphs
    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)

    colour_val = ACCENT_HEX if is_light_bg else "FFFFFF"

    # Paragraph 1: Intro text
    intro_p = etree.SubElement(txBody, qn('a:p'))
    intro_pPr = etree.SubElement(intro_p, qn('a:pPr'))
    intro_lnSpc = etree.SubElement(intro_pPr, qn('a:lnSpc'))
    etree.SubElement(intro_lnSpc, qn('a:spcPct'), val="120000")

    intro_r = etree.SubElement(intro_p, qn('a:r'))
    intro_rPr = etree.SubElement(intro_r, qn('a:rPr'), sz="2200", dirty="0")
    intro_fill = etree.SubElement(intro_rPr, qn('a:solidFill'))
    etree.SubElement(intro_fill, qn('a:srgbClr'), val=colour_val)
    etree.SubElement(intro_rPr, qn('a:latin'), typeface=FONT_FAMILY)
    intro_t = etree.SubElement(intro_r, qn('a:t'))
    intro_t.text = intro_text

    # Paragraph 2: Empty line for spacing
    empty_p = etree.SubElement(txBody, qn('a:p'))
    empty_pPr = etree.SubElement(empty_p, qn('a:pPr'))
    empty_lnSpc = etree.SubElement(empty_pPr, qn('a:lnSpc'))
    etree.SubElement(empty_lnSpc, qn('a:spcPct'), val="120000")
    etree.SubElement(empty_p, qn('a:endParaRPr'), dirty="0")

    # Bullet paragraphs with arrow (Wingdings Ø)
    for bullet_text in bullets:
        bullet_p = etree.SubElement(txBody, qn('a:p'))
        bullet_pPr = etree.SubElement(bullet_p, qn('a:pPr'), marL="342900", indent="-342900")
        bullet_lnSpc = etree.SubElement(bullet_pPr, qn('a:lnSpc'))
        etree.SubElement(bullet_lnSpc, qn('a:spcPct'), val="120000")
        # Arrow bullet using Wingdings
        etree.SubElement(bullet_pPr, qn('a:buFont'), typeface="Wingdings", pitchFamily="2", charset="2")
        etree.SubElement(bullet_pPr, qn('a:buChar'), char="\u00D8")  # Ø = arrow in Wingdings

        bullet_r = etree.SubElement(bullet_p, qn('a:r'))
        bullet_rPr = etree.SubElement(bullet_r, qn('a:rPr'), sz="2200", dirty="0")
        bullet_fill = etree.SubElement(bullet_rPr, qn('a:solidFill'))
        etree.SubElement(bullet_fill, qn('a:srgbClr'), val=colour_val)
        etree.SubElement(bullet_rPr, qn('a:latin'), typeface=FONT_FAMILY)
        bullet_t = etree.SubElement(bullet_r, qn('a:t'))
        bullet_t.text = bullet_text


# =============================================================================
# SLIDE MANAGEMENT
# =============================================================================

def add_slide(prs: Presentation, layout_key: str):
    """Add slide using layout key from LAYOUTS dictionary."""
    return prs.slides.add_slide(prs.slide_layouts[LAYOUTS[layout_key]])


# =============================================================================
# ANIMATIONS
# =============================================================================

def add_dissolve_animations(slide):
    """
    Add 0.5s dissolve on-click animation to all text elements.

    Animation behaviour:
    - Title: animates as single unit
    - Subtitle: animates paragraph-by-paragraph
    - Body: animates paragraph-by-paragraph (intro, then each bullet)

    Do NOT call this for fixed slides (About, CTA) — they have
    pre-set animations in the template.
    """
    shapes_info = []

    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        spid = shape.shape_id
        para_count = len(shape.text_frame.paragraphs)

        # Check if this should animate by paragraph
        is_body = False
        try:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                idx = shape.placeholder_format.idx
                if idx >= 10 or idx == 1:  # Body or subtitle
                    is_body = True
        except:
            pass

        shapes_info.append((spid, para_count, is_body))

    if not shapes_info:
        return

    # Build animation XML
    ctn_id = 1
    child_pars = []
    bld_items = []

    for spid, para_count, is_body in shapes_info:
        if is_body and para_count > 1:
            for para_idx in range(para_count):
                ctn_id += 1
                base_id = ctn_id

                par_xml = f'''
                <p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
                    <p:cTn id="{base_id}" fill="hold">
                        <p:stCondLst>
                            <p:cond delay="indefinite"/>
                        </p:stCondLst>
                        <p:childTnLst>
                            <p:par>
                                <p:cTn id="{base_id + 1}" fill="hold">
                                    <p:stCondLst>
                                        <p:cond delay="0"/>
                                    </p:stCondLst>
                                    <p:childTnLst>
                                        <p:par>
                                            <p:cTn id="{base_id + 2}" presetID="9" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="clickEffect">
                                                <p:stCondLst>
                                                    <p:cond delay="0"/>
                                                </p:stCondLst>
                                                <p:childTnLst>
                                                    <p:set>
                                                        <p:cBhvr>
                                                            <p:cTn id="{base_id + 3}" dur="1" fill="hold">
                                                                <p:stCondLst>
                                                                    <p:cond delay="0"/>
                                                                </p:stCondLst>
                                                            </p:cTn>
                                                            <p:tgtEl>
                                                                <p:spTgt spid="{spid}">
                                                                    <p:txEl>
                                                                        <p:pRg st="{para_idx}" end="{para_idx}"/>
                                                                    </p:txEl>
                                                                </p:spTgt>
                                                            </p:tgtEl>
                                                            <p:attrNameLst>
                                                                <p:attrName>style.visibility</p:attrName>
                                                            </p:attrNameLst>
                                                        </p:cBhvr>
                                                        <p:to>
                                                            <p:strVal val="visible"/>
                                                        </p:to>
                                                    </p:set>
                                                    <p:animEffect transition="in" filter="dissolve">
                                                        <p:cBhvr>
                                                            <p:cTn id="{base_id + 4}" dur="500"/>
                                                            <p:tgtEl>
                                                                <p:spTgt spid="{spid}">
                                                                    <p:txEl>
                                                                        <p:pRg st="{para_idx}" end="{para_idx}"/>
                                                                    </p:txEl>
                                                                </p:spTgt>
                                                            </p:tgtEl>
                                                        </p:cBhvr>
                                                    </p:animEffect>
                                                </p:childTnLst>
                                            </p:cTn>
                                        </p:par>
                                    </p:childTnLst>
                                </p:cTn>
                            </p:par>
                        </p:childTnLst>
                    </p:cTn>
                </p:par>
                '''
                child_pars.append(par_xml)
                ctn_id += 4

            bld_items.append(f'<p:bldP spid="{spid}" grpId="0" build="p"/>')
        else:
            ctn_id += 1
            base_id = ctn_id

            par_xml = f'''
            <p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
                <p:cTn id="{base_id}" fill="hold">
                    <p:stCondLst>
                        <p:cond delay="indefinite"/>
                    </p:stCondLst>
                    <p:childTnLst>
                        <p:par>
                            <p:cTn id="{base_id + 1}" fill="hold">
                                <p:stCondLst>
                                    <p:cond delay="0"/>
                                </p:stCondLst>
                                <p:childTnLst>
                                    <p:par>
                                        <p:cTn id="{base_id + 2}" presetID="9" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="clickEffect">
                                            <p:stCondLst>
                                                <p:cond delay="0"/>
                                            </p:stCondLst>
                                            <p:childTnLst>
                                                <p:set>
                                                    <p:cBhvr>
                                                        <p:cTn id="{base_id + 3}" dur="1" fill="hold">
                                                            <p:stCondLst>
                                                                <p:cond delay="0"/>
                                                            </p:stCondLst>
                                                        </p:cTn>
                                                        <p:tgtEl>
                                                            <p:spTgt spid="{spid}"/>
                                                        </p:tgtEl>
                                                        <p:attrNameLst>
                                                            <p:attrName>style.visibility</p:attrName>
                                                        </p:attrNameLst>
                                                    </p:cBhvr>
                                                    <p:to>
                                                        <p:strVal val="visible"/>
                                                    </p:to>
                                                </p:set>
                                                <p:animEffect transition="in" filter="dissolve">
                                                    <p:cBhvr>
                                                        <p:cTn id="{base_id + 4}" dur="500"/>
                                                        <p:tgtEl>
                                                            <p:spTgt spid="{spid}"/>
                                                        </p:tgtEl>
                                                    </p:cBhvr>
                                                </p:animEffect>
                                            </p:childTnLst>
                                        </p:cTn>
                                    </p:par>
                                </p:childTnLst>
                            </p:cTn>
                        </p:par>
                    </p:childTnLst>
                </p:cTn>
            </p:par>
            '''
            child_pars.append(par_xml)
            ctn_id += 4
            bld_items.append(f'<p:bldP spid="{spid}" grpId="0"/>')

    all_pars = ''.join(child_pars)
    all_blds = ''.join(bld_items)

    timing_xml = f'''
    <p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
        <p:tnLst>
            <p:par>
                <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
                    <p:childTnLst>
                        <p:seq concurrent="1" nextAc="seek">
                            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                                <p:childTnLst>
                                    {all_pars}
                                </p:childTnLst>
                            </p:cTn>
                            <p:prevCondLst>
                                <p:cond evt="onPrev" delay="0">
                                    <p:tgtEl>
                                        <p:sldTgt/>
                                    </p:tgtEl>
                                </p:cond>
                            </p:prevCondLst>
                            <p:nextCondLst>
                                <p:cond evt="onNext" delay="0">
                                    <p:tgtEl>
                                        <p:sldTgt/>
                                    </p:tgtEl>
                                </p:cond>
                            </p:nextCondLst>
                        </p:seq>
                    </p:childTnLst>
                </p:cTn>
            </p:par>
        </p:tnLst>
        <p:bldLst>
            {all_blds}
        </p:bldLst>
    </p:timing>
    '''

    try:
        timing_elm = etree.fromstring(timing_xml)
        existing = slide._element.find(qn('p:timing'))
        if existing is not None:
            slide._element.remove(existing)
        slide._element.append(timing_elm)
    except Exception as e:
        print(f"  Warning: Animation error: {e}")


# =============================================================================
# EXAMPLE OUTLINE — REPLACE WITH YOUR CONTENT
# =============================================================================
#
# The script generates slides dynamically from this outline.
# Modify the OUTLINE dict to change presentation content.
#
# Supported slide types:
#   - "content": Title + Subtitle + Body (intro paragraph + bullets)
#   - "quote": Quote slide
#
# Section options:
#   - "section_type": "blue" (default), "pale", or "none" (skip section header)
#
# Layout alternation: Content slides alternate between white/pale automatically
# unless you specify "layout" explicitly in the slide dict.

OUTLINE = {
    "title": "PRESENTATION TITLE",
    "subtitle": "Your Subtitle Here",
    "thank_you_subtitle": "Contact information or next steps",

    "sections": [
        {
            "name": "SECTION ONE",
            "subtitle": "Section description",
            "slides": [
                {
                    "type": "content",
                    "title": "SLIDE TITLE",
                    "subtitle": "Short Subtitle (5-10 words)",
                    "intro": "This is the intro paragraph explaining the context. Keep it to 1-2 sentences that frame the bullet points.",
                    "bullets": [
                        "First key point",
                        "Second key point",
                        "Third key point",
                        "Fourth key point (maximum)"
                    ]
                }
            ]
        },
        {
            "name": "SECTION TWO",
            "subtitle": "Another section",
            "slides": [
                {
                    "type": "content",
                    "title": "ANOTHER SLIDE",
                    "subtitle": "More Content Here",
                    "intro": "Another intro paragraph providing context for the following points.",
                    "bullets": [
                        "Point A",
                        "Point B",
                        "Point C"
                    ]
                },
                {
                    "type": "content",
                    "title": "THIRD SLIDE",
                    "subtitle": "Even More Content",
                    "intro": "Sections can have multiple content slides. The script handles this automatically.",
                    "bullets": [
                        "First item",
                        "Second item"
                    ]
                }
            ]
        },
        {
            "name": "LIMITATIONS",
            "subtitle": "What to watch out for",
            "section_type": "pale",  # Use pale section for caveats
            "slides": [
                {
                    "type": "content",
                    "title": "CURRENT CONSTRAINTS",
                    "subtitle": "Room for Improvement",
                    "intro": "Every solution has limitations. Be transparent about them.",
                    "bullets": [
                        "Limitation one",
                        "Limitation two",
                        "Limitation three"
                    ]
                }
            ]
        },
        {
            "name": "Q&A",
            "subtitle": "Questions?",
            "slides": []  # Section header only, no content slides
        }
    ]
}


# =============================================================================
# PRESENTATION CREATION
# =============================================================================

def create_presentation(outline: dict = None, output_path: Path = None):
    """
    Create presentation from structured outline.

    Args:
        outline: Presentation outline dict. Uses OUTLINE constant if not provided.
        output_path: Where to save the output. Uses OUTPUT_PATH if not provided.
    """
    if outline is None:
        outline = OUTLINE
    if output_path is None:
        output_path = OUTPUT_PATH

    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE_PATH}")

    # Convert template
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = Path(tmp.name)
    convert_potx_to_pptx(TEMPLATE_PATH, tmp_path)

    try:
        prs = Presentation(str(tmp_path))

        # Remove all existing slides
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        slide_num = 0
        content_slide_count = 0  # For alternating white/pale

        # === FIXED OPENING SEQUENCE ===

        # Title slide
        slide_num += 1
        slide = add_slide(prs, "title")
        set_text_simple(slide.placeholders[0], outline["title"])
        set_text_simple(slide.placeholders[1], outline["subtitle"])
        add_dissolve_animations(slide)
        print(f"{slide_num}. Title slide")

        # Agenda slide — auto-generate from section names
        slide_num += 1
        slide = add_slide(prs, "menu")
        set_text_simple(slide.placeholders[0], "AGENDA", is_light_bg=True)
        agenda_items = [s["name"] for s in outline["sections"] if s.get("section_type") != "none"]
        agenda_text = "\n".join(agenda_items)
        set_text_simple(slide.placeholders[1], agenda_text, is_light_bg=True)
        add_dissolve_animations(slide)
        print(f"{slide_num}. Agenda slide")

        # About slide (Fixed) — comment out if your template doesn't have one
        slide_num += 1
        slide = add_slide(prs, "about")
        print(f"{slide_num}. About slide (fixed)")

        # === SECTIONS ===

        for section in outline["sections"]:
            section_type = section.get("section_type", "blue")

            # Section header (unless section_type is "none")
            if section_type != "none":
                slide_num += 1
                if section_type == "pale":
                    slide = add_slide(prs, "section_pale")
                    set_text_simple(slide.placeholders[0], section["name"], is_light_bg=True)
                    set_text_simple(slide.placeholders[1], section["subtitle"], is_light_bg=True)
                else:
                    slide = add_slide(prs, "section")
                    set_text_simple(slide.placeholders[0], section["name"])
                    set_text_simple(slide.placeholders[1], section["subtitle"])
                add_dissolve_animations(slide)
                print(f"{slide_num}. Section: {section['name']}")

            # Content slides
            for content in section.get("slides", []):
                slide_num += 1

                if content["type"] == "content":
                    # Alternate between white and pale
                    layout_key = content.get("layout")
                    if layout_key is None:
                        layout_key = "content_white" if content_slide_count % 2 == 0 else "content_pale"
                    content_slide_count += 1

                    slide = add_slide(prs, layout_key)
                    set_text_simple(slide.placeholders[0], content["title"], is_light_bg=True)
                    set_text_simple(slide.placeholders[1], content["subtitle"], is_light_bg=True)
                    set_body_with_bullets(
                        slide.placeholders[BODY_PLACEHOLDER_IDX],
                        content["intro"],
                        content["bullets"],
                        is_light_bg=True
                    )
                    add_dissolve_animations(slide)
                    print(f"{slide_num}. {content['title']}")

                elif content["type"] == "quote":
                    slide = add_slide(prs, "quote")
                    set_text_simple(slide.placeholders[0], content["quote"], is_light_bg=True)
                    set_text_simple(slide.placeholders[1], content.get("attribution", ""), is_light_bg=True)
                    add_dissolve_animations(slide)
                    print(f"{slide_num}. Quote")

        # === FIXED CLOSING SEQUENCE ===

        # CTA slide (Fixed) — comment out if your template doesn't have one
        slide_num += 1
        slide = add_slide(prs, "cta")
        print(f"{slide_num}. CTA (fixed)")

        # Thank You slide
        slide_num += 1
        slide = add_slide(prs, "thank_you")
        set_text_simple(slide.placeholders[0], "THANK YOU")
        set_text_simple(slide.placeholders[1], outline.get("thank_you_subtitle", ""))
        add_dissolve_animations(slide)
        print(f"{slide_num}. Thank You")

        # Save
        prs.save(str(output_path))
        print(f"\nSaved: {output_path}")
        print(f"Total slides: {len(prs.slides)}")

    finally:
        if tmp_path.exists():
            tmp_path.unlink()


# =============================================================================
# ENTRY POINT
# =============================================================================

if __name__ == "__main__":
    create_presentation()
